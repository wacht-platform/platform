from pptx import Presentation
from pptx.util import Inches


def build_image_showcase(output_path: str, image_path: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Architecture Overview"
    slide.shapes.add_picture(image_path, Inches(1.0), Inches(1.6), width=Inches(7.5))
    note = slide.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(7.5), Inches(0.5))
    note.text_frame.text = "Current system boundaries and key data paths"
    prs.save(output_path)


if __name__ == "__main__":
    build_image_showcase("/tmp/architecture-showcase.pptx", "/tmp/diagram.png")
