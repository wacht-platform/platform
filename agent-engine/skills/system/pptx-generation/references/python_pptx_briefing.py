from pptx import Presentation
from pptx.util import Inches, Pt


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)


def build_deck(output_path: str) -> None:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly Strategy Brief"
    title_slide.placeholders[1].text = "Executive summary and recommendations"

    add_bullets_slide(
        prs,
        "What Changed",
        [
            "Buyer expectations shifted toward measurable ROI.",
            "Competitive pricing pressure increased in SMB segments.",
            "Enterprise expansion now depends on trust and governance signals.",
        ],
    )

    add_bullets_slide(
        prs,
        "Recommended Moves",
        [
            "Package premium reporting into the enterprise tier.",
            "Sharpen positioning around operational trust.",
            "Turn repeatable workflows into template-driven onboarding.",
        ],
    )

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "Competitor Snapshot"
    table = table_slide.shapes.add_table(4, 3, Inches(0.7), Inches(1.8), Inches(8.0), Inches(2.5)).table
    table.cell(0, 0).text = "Vendor"
    table.cell(0, 1).text = "Strength"
    table.cell(0, 2).text = "Risk"
    rows = [
        ("InboxDoctor", "Workflow depth", "Lower awareness"),
        ("Lemlist", "Top-of-funnel reach", "Less defensible"),
        ("ZeroBounce", "Category trust", "Broader but less specific"),
    ]
    for row_index, values in enumerate(rows, start=1):
        for col_index, value in enumerate(values):
            table.cell(row_index, col_index).text = value

    prs.save(output_path)


if __name__ == "__main__":
    build_deck("/tmp/example_briefing.pptx")
