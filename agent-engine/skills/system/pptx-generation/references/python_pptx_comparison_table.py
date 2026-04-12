from pptx import Presentation
from pptx.util import Inches


def build_comparison_deck(output_path: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Option Comparison"

    table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.5), Inches(8.5), Inches(2.5)).table
    data = [
        ["Option", "Cost", "Speed", "Risk"],
        ["A", "$24k", "Fast", "Medium"],
        ["B", "$31k", "Medium", "Low"],
        ["C", "$19k", "Fast", "High"],
    ]
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    prs.save(output_path)


if __name__ == "__main__":
    build_comparison_deck("/tmp/option-comparison.pptx")
