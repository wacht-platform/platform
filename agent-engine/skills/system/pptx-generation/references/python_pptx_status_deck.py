from pptx import Presentation
from pptx.util import Inches


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.text = bullets[0]
    for bullet in bullets[1:]:
        frame.add_paragraph().text = bullet


def build_status_deck(output_path: str) -> None:
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "Project Status Review"
    title.placeholders[1].text = "Progress, risks, and next actions"

    add_bullet_slide(prs, "Completed", ["Backend migration complete", "Monitoring added", "Docs refreshed"])
    add_bullet_slide(prs, "Current Risks", ["Vendor delay remains possible", "Testing capacity is limited"])
    add_bullet_slide(prs, "Next Actions", ["Finish regression run", "Decide cutover date", "Prepare stakeholder update"])

    prs.save(output_path)


if __name__ == "__main__":
    build_status_deck("/tmp/project-status.pptx")
